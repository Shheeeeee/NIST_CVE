import requests
import datetime
from urllib.parse import quote
import time    
from bs4 import BeautifulSoup
# import re

import json

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches  # Import the Inches function
import copy

# from googletrans import Translator

from itertools import cycle
from shutil import get_terminal_size
from threading import Thread
from time import sleep
from sys import stdout

#################################################################################
#                                 ANIMATION
#################################################################################

class Loader: # pour faire joli quand on le lance manuellement
    def __init__(self, desc="Tentative de connexion ", end="", timeout=0.1):
        self.desc = desc
        self.end = end
        self.timeout = timeout

        self._thread = Thread(target=self._animate, daemon=True)
        self.steps = ["⢿ ", "⣻ ", "⣽ ", "⣾ ", "⣷ ", "⣯ ", "⣟ ", "⡿ "]
        self.done = False

    def start(self):
        self._thread.start()
        return self

    def _animate(self):
        for c in cycle(self.steps):
            if self.done:
                break
            text = f"{self.desc}{self.attempt} {c}"
            print(f"\r{text}", flush=True, end="")
            stdout.flush()
            time.sleep(self.timeout)

    def __enter__(self):
        self.start()

    def stop(self):
        self.done = True
        cols = get_terminal_size((80, 20)).columns
        print("\r" + " " * cols, end="", flush=True)
        print(f"\r{self.end}", flush=True)

    def set_attempt(self, attempt):
        self.attempt = attempt

    def __exit__(self, exc_type, exc_value, tb):
        self.stop()

#################################################################################
#                                 API CALL
#################################################################################

def make_api_request(url, max_attempts=5): # CALL API NIST (url à renseigner en arg)
    loader = Loader()
    loader_thread = Thread(target=loader.start, daemon=True)
    loader_thread.start()

    try:
        for attempt in range(1, max_attempts + 1):
            loader.set_attempt(attempt)
            #print(f"Connexion à l'API en cours...", end="")

            response = requests.get(url)

            if response.status_code == 200:
                data = response.json()
                loader.stop()
                print("\nConnexion réussie !")
                return data

            time.sleep(6)
    finally:
        loader.stop()
        if attempt == max_attempts:
            print(f"Impossible de se connecter à l'API après {max_attempts} tentatives. Veuillez réessayer plus tard.")
            exit()

def scrape_composants(cve_id): # scraping pour récupérer le produit et les versions
    try:
        url = f"https://cveawg.mitre.org/api/cve/{cve_id}"

        response = requests.get(url)
        response.raise_for_status()
        
        # Charger le contenu JSON
        data = json.loads(response.text)
        
        # Accéder à la clé "containers" pour obtenir les informations sur les produits affectés
        containers = data.get('containers', {}).get('cna', {}).get('affected', [])
        
        # Créer des ensembles (set) vides pour stocker temporairement les informations uniques sur les fournisseurs
        unique_vendors = set()
        products = []
        all_versions = []
        
        for container in containers:
            vendor = container.get('vendor')
            product = container.get('product')
            versions = container.get('versions')
            
            # Ajouter le fournisseur dans l'ensemble des fournisseurs uniques
            unique_vendors.add(vendor)
            
            # Enregistrer les informations de chaque produit
            products.append(product)
            all_versions.extend(versions)
        
        # Convertir l'ensemble des fournisseurs uniques en liste
        vendors = list(unique_vendors)
        
        return vendors, products, all_versions

    except Exception as e:
        print(f"Une erreur s'est produite : {e}")
        return None, None, None

def affichage_products(products, versions): # Mise en forme des versions et produits
    composants = ""
    try:
        for i in range(len(products)):
            version_info = versions[i]
            less_than_version = version_info.get('lessThan')
            version = version_info.get('version')
            
            # Utiliser 'lessThan' s'il existe, sinon utiliser 'version'
            version_to_use = less_than_version if less_than_version is not None else version
            
            composants = composants + f"{products[i]} anterieur à {version_to_use}, "
        composants = composants[:-2]
        return composants

    except Exception as e:
        print(f"Une erreur s'est produite : {e}")
        return None


#################################################################################
#                                    CVE
#################################################################################

def CVE(start_date, end_date):

    # URL de l'API du NIST pour les CVE avec les dates de début et de fin spécifiées
    # url = f'https://services.nvd.nist.gov/rest/json/cves/2.0/?lastModStartDate={encoded_start_date}&lastModEndDate={encoded_end_date}'
    url = f'https://services.nvd.nist.gov/rest/json/cves/2.0/?pubStartDate={start_date}&pubEndDate={end_date}'

    data = make_api_request(url)

    if len(data) != 0:
        cve_items = data['vulnerabilities']
        cve_list = []
        
        # Parcourir chaque CVE et extraire les informations
        for cve_item in cve_items:

            vuln_status = cve_item['cve']['vulnStatus']
            cve_id = cve_item['cve']['id']

            severity = None
            for i in range(1, 35): # je recupere la severity de la CVE
                key = f'cvssMetricV{i}'
                if key in cve_item['cve']['metrics']:
                    cve_metrics = cve_item['cve']['metrics'][key][0]

                    #################
                    # severity
                    #################

                    try:
                        severity = cve_metrics['baseScore']
                    except KeyError:
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                            severity = cve_metrics['cvssData']['baseScore']
                        except KeyError:
                            # Si 'attackVector' n'existe pas, définir attack_vector comme None
                            severity = None
                        continue

            #################
            # Tri des CVE, ne prends que les CVE avec plus de 8 de score et des status autres que Rejected et modified (qui sont les CVE non acceptés ou doublons par rapport à une ancienne)
            #################

            # ici 2 solutions soit on a une severity, soit pas, donc on va récupérer toutes les CVE
            # qui ont une severity au dessus de 8 ou qui est NONE

            # entre 7 et 8 et pas les None
            # if (severity is None) or (severity is not None and severity < 7 or severity > 8) or vuln_status == "Rejected" or vuln_status == "Modified":

            if (severity is not None and severity < 8) or vuln_status == "Rejected" or vuln_status == "Modified":
                print(f"Skipping CVE {cve_id} due to low severity or rejected/modified status")
            else:
                #################
                # Récupération des inforamtions de base
                #################
     
                published = cve_item['cve']['published']
                lastModified = cve_item['cve']['lastModified']
                description = cve_item['cve']['descriptions'][0]['value']
                
                #################
                # source    
                #################
                try:
                    source = cve_item['cve']['references'][0]['url']
                except Exception as e:
                    source = None

                try:
                    source2 = cve_item['cve']['references'][1]['url']
                except Exception as e:
                    source2 = None
                
                try:
                    # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]
                    source3 = cve_item['cve']['references'][2]['url']
                except Exception as e:
                    source3 = None
                
                #produit = scrape_website(cve_id) Scrap du nist directement (pas dingue)
                
                #scrpe des produits sur CVE.com
                produit, products, versions = scrape_composants(cve_id)
                composants = affichage_products(products, versions)

                #initialisation des variables et elles se reinitialisent entre chaque CVE

                vector_string = None
                attack_vector = None
                attack_complexity = None
                privileges_required = None
                user_interaction = None
                scope = None
                confidentiality_impact = None
                integrity_impact = None
                availability_impact = None
                base_severity = None

                for i in range(1, 35): # Essayer différentes clés pour extraire le baseScore
                    key = f'cvssMetricV{i}'  # Construire la clé correspondante
                    if key in cve_item['cve']['metrics']:
                        cve_metrics = cve_item['cve']['metrics'][key][0]

                        # TODO UN BOUCLE PARCEQUE LA C'EST IMMONDE EN DESSOUS
                        
                        #################
                        # vectorString
                        #################
                        try:
                            # Essayer d'accéder à 'vectorString' dans cve_item['cve']['metrics'][key][0]['cvssData']
                            vector_string = cve_metrics['vectorString']
                        except Exception as e:
                            try:
                                # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                                vector_string = cve_metrics['cvssData']['vectorString']
                            except Exception as e:
                                # Si 'attackVector' n'existe pas, définir attack_vector comme None
                                vector_string = None

                        #################
                        # attackVector
                        #################
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]
                            attack_vector = cve_metrics['attackVector']
                        except KeyError:
                            try:
                                # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                                attack_vector = cve_metrics['cvssData']['attackVector']
                            except KeyError:
                                # Si 'attackVector' n'existe pas, définir attack_vector comme None
                                attack_vector = None
                        
                        #################
                        # attack_complexity 
                        #################
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]
                            attack_complexity = cve_metrics['attackComplexity']
                        except KeyError:
                            try:
                                # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                                attack_complexity = cve_metrics['cvssData']['attackComplexity']
                            except KeyError:
                                # Si 'attackVector' n'existe pas, définir attack_vector comme None
                                attack_complexity = None
                        
                        #################
                        # privileges_required  
                        #################
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]
                            privileges_required = cve_metrics['privilegesRequired']
                        except KeyError:
                            try:
                                # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                                privileges_required = cve_metrics['cvssData']['privilegesRequired']
                            except KeyError:
                                # Si 'attackVector' n'existe pas, définir attack_vector comme None
                                privileges_required = None
                        
                        #################
                        # userInteraction  
                        #################
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]
                            user_interaction = cve_metrics['userInteraction']
                        except KeyError:
                            try:
                                # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                                user_interaction = cve_metrics['cvssData']['userInteraction']
                            except KeyError:
                                # Si 'attackVector' n'existe pas, définir attack_vector comme None
                                user_interaction = None
                        
                        #################
                        # scope   
                        #################
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]
                            scope = cve_metrics['scope']
                        except KeyError:
                            try:
                                # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                                scope = cve_metrics['cvssData']['scope']
                            except KeyError:
                                # Si 'attackVector' n'existe pas, définir attack_vector comme None
                                scope = None

                        #################
                        # confidentiality_impact   
                        #################
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]
                            confidentiality_impact = cve_metrics['confidentialityImpact']
                        except KeyError:
                            try:
                                # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                                confidentiality_impact = cve_metrics['cvssData']['confidentialityImpact']
                            except KeyError:
                                # Si 'attackVector' n'existe pas, définir attack_vector comme None
                                confidentiality_impact = None

                        #################
                        # integrityImpact   
                        #################
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]
                            integrity_impact = cve_metrics['integrityImpact']
                        except KeyError:
                            try:
                                # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                                integrity_impact = cve_metrics['cvssData']['integrityImpact']
                            except KeyError:
                                # Si 'attackVector' n'existe pas, définir attack_vector comme None
                                integrity_impact = None
                        
                        #################
                        # availability_impact    
                        #################
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]
                            availability_impact = cve_metrics['availabilityImpact']
                        except KeyError:
                            try:
                                # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                                availability_impact = cve_metrics['cvssData']['availabilityImpact']
                            except KeyError:
                                # Si 'attackVector' n'existe pas, définir attack_vector comme None
                                availability_impact = None

                        #################
                        # baseSeverity    
                        #################
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]
                            base_severity = cve_metrics['baseSeverity']
                        except KeyError:
                            try:
                                # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                                base_severity = cve_metrics['cvssData']['baseSeverity']
                            except KeyError:
                                # Si 'attackVector' n'existe pas, définir attack_vector comme None
                                base_severity = None
                    
                # mise en forme rapide a changer c'est vraiment pas beau
                try: # parfois produit est None alors qu'il devrait pas ce grand fou mais osef on le rempli si il est pas content
                    if produit[0] == 'n/a':
                        produit[0] = ""
                except Exception as e:
                    if produit == None:
                        produit = ['']  
                
                if composants == 'None':
                    composants = ""

                if composants == 'n/a anterieur à n/a':
                    composants = ""

                # Remplir le tableau cve_info
                cve_info = {
                'produit': produit[0],
                'cve_id': cve_id,
                'published': published,
                'lastModified': lastModified,
                'vector_string': vector_string,
                'attack_vector': attack_vector,
                'attack_complexity': attack_complexity,
                'privileges_required': privileges_required,
                'user_interaction': user_interaction,
                'scope': scope,
                'confidentiality_impact': confidentiality_impact,
                'integrity_impact': integrity_impact,
                'availability_impact': availability_impact,
                'severity': severity,
                'base_severity': base_severity,
                'descriptions': description,
                'source' : source,
                'source2' : source2,
                'source3' : source3,
                'composants' : composants

                }

                #ajouter la ligne au tableau
                cve_list.append(cve_info)

        # on print (dans un soucis de debug uniquement c'est inutile autrement)                          
        for cve_info in cve_list:

            print(f"CVE ID: {cve_info['cve_id']}")
            print(f"produit: {cve_info['produit']}")
            print(f"publiée : {cve_info['published']}  modifiée : {cve_info['lastModified']}")

            print(f"Vuln Status: {vuln_status}")
            print(f"Severity: {cve_info['severity']}")
            print(f"base_severity: {cve_info['base_severity']}")

            print(f"Vector String: {cve_info['vector_string']}")
            print(f"Attack Vector: {cve_info['attack_vector']}")
            print(f"attack_complexity: {cve_info['attack_complexity']}")
            print(f"privileges_required: {cve_info['privileges_required']}")
            print(f"user_interaction: {cve_info['user_interaction']}")
            print(f"scope: {cve_info['scope']}")
            print(f"confidentiality_impact: {cve_info['confidentiality_impact']}")
            print(f"integrity_impact: {cve_info['integrity_impact']}")
            print(f"availability_impact: {cve_info['availability_impact']}")

            print(f"Descriptions: {cve_info['descriptions']}")
            print(f"Source: {cve_info['source']}")
            print()
    
    else:
        print(f'Aucune CVE trouvée pour la plage de dates spécifiée.')
    
    return cve_list
    
#################################################################################
#                                 POWERPOINT
#################################################################################

def powerpoint(cve_list, start_date, end_date):
    
    # Le nom du template (modèle) et du fichier de sortie
    template_path = "Bulletin_de_veille_TEMPLATE.pptx"

    # Ouvrir le modèle
    presentation = Presentation(template_path)

    # Dupliquer la deuxième diapo (index 1) quatre fois
    for _ in range(len(cve_list) - 1):
        duplicate_slide(presentation, 1)
    
    diapo_index = 1  # Commencer à la deuxième diapo (index 1)
    
    for i in range(len(cve_list)):  # Modifier jusqu'à la cinquième diapo (index 4)
        slide = presentation.slides[diapo_index]

        cve_list[diapo_index - 1]['source2'], cve_list[diapo_index - 1]['source3'] = source_forme(cve_list[diapo_index - 1]['source2'], cve_list[diapo_index - 1]['source3'])
        
        cell_coords = [(2, 2), (3, 2), (4, 2), (5, 2), (2, 5), (3, 5), (4, 5), (5, 5), (7, 1), (10, 1), (6,1)] 
        cell_values = [cve_list[diapo_index - 1]['attack_vector'], cve_list[diapo_index - 1]['attack_complexity'], cve_list[diapo_index - 1]['privileges_required'], cve_list[diapo_index - 1]['user_interaction'], cve_list[diapo_index - 1]['scope'], cve_list[diapo_index - 1]['confidentiality_impact'], cve_list[diapo_index - 1]['integrity_impact'], cve_list[diapo_index - 1]['availability_impact'], cve_list[diapo_index - 1]['descriptions'], f"{str(cve_list[diapo_index - 1]['source'])}\n{cve_list[diapo_index - 1]['source2']}\n{cve_list[diapo_index - 1]['source3']}", cve_list[diapo_index - 1]['composants']]  

        for i in range(len(cell_coords)):
            modify_table_cell_black(slide, cell_coords[i], cell_values[i])
            
        cell_coords = [(0, 1), (2, 7) , (9, 7)]  
        cell_values = [cve_list[diapo_index - 1]['vector_string'], f"{str(cve_list[diapo_index - 1]['severity'])}\n{cve_list[diapo_index - 1]['base_severity']}", cve_list[diapo_index - 1]['base_severity']]  

        for i in range(len(cell_coords)):
            modify_table_cell_white(slide, cell_coords[i], cell_values[i])
        
        titre(slide, f"{cve_list[diapo_index - 1]['cve_id']} {cve_list[diapo_index - 1]['produit']}")

        diapo_index += 1

    # Enregistrer la présentation modifiée dans un fichier
    nom_fichier = "CERT - Bulletin_de_veille_" + end_date + ".pptx"
    presentation.save(nom_fichier)
    print(f"La présentation a été sauvegardée dans '{nom_fichier}' avec succès.")

def source_forme(source2, source3):
    if source2 == None:
        source2 = ""
    if source3 == None:
        source3 = ""
    return source2, source3

def titre(slide, titre):
    title_shape = slide.shapes.title

    # Modifiez le texte du titre

    title_shape.text = titre

    # Définir le style du texte du titre
    font = title_shape.text_frame.paragraphs[0].font
    font.name = "Poppins SemiBold"
    #font.bold = True
    font.size = Pt(36)

    left = Pt(55)  # Définir la position horizontale (gauche) en points
    top = Pt(90)   # Définir la position verticale (haut) en points
    title_shape.left = left
    title_shape.top = top

def duplicate_slide(pres, index):
    try:
        template = pres.slides[index]
        blank_slide_layout = template.slide_layout
    except:
        blank_slide_layout = pres.slide_layouts[0]

    copied_slide = pres.slides.add_slide(blank_slide_layout)

    for shp in template.shapes:
        if not shp.has_text_frame or shp.text_frame.text != "":  # Vérifier si la forme contient du texte
            el = shp.element
            newel = copy.deepcopy(el)
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return copied_slide


def modify_table_cell_black(slide, cell_coords, cell_value):
    row, col = cell_coords
    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break

    if table is None:
        print("Aucun tableau trouvé dans la diapositive.")
        return

    # Modifier le texte dans la cellule de la ligne et colonne spécifiées
    cell = table.cell(row, col)
    text_frame = cell.text_frame
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()

    # Modifier la mise en forme du texte
    run.text = str(cell_value)
    run.font.size = Pt(11)
    run.font.name = "Poppins"
    run.font.color.rgb = RGBColor(0, 0, 0)  # Noir


def modify_table_cell_white(slide, cell_coords, cell_value):
    row, col = cell_coords

    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break

    if table is None:
        print("Aucun tableau trouvé dans la diapositive.")
        return

    # Modifier le texte dans la cellule de la ligne et colonne spécifiées
    cell = table.cell(row, col)
    text_frame = cell.text_frame
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()

    # Modifier la mise en forme du texte
    run.text = str(cell_value)
    run.font.size = Pt(11)
    run.font.name = "Poppins"
    run.font.bold = True  # Pour mettre le texte en gras
    run.font.color.rgb = RGBColor(255, 255, 255)  # Pour mettre le texte en blanc (255, 255, 255 correspond au blanc en RGB)

#################################################################################
#                              MISE EN FORME
#################################################################################

def trad_vectors(cve_list):
    # oupsi j'ai glissé, c'est vraiment pas beau tout ça, grosse dédicasse à la personne qui lit ce code
    # je suis désolé, partie à revoir
    for cve_info in cve_list:
        if cve_info['attack_vector'] == 'NETWORK':
            cve_info['attack_vector'] = "Réseau"
        elif cve_info['attack_vector'] == 'LOCAL':
            cve_info['attack_vector'] = "Local"
        elif cve_info['attack_vector'] == 'PHYSICAL':
            cve_info['attack_vector'] = "Réseau"
        elif cve_info['attack_vector'] == 'ADJACENT NETWORK':
            cve_info['attack_vector'] = "Réseau adjacent"

        if cve_info['base_severity'] == 'HIGH':
            cve_info['base_severity'] = "Élevée"

        if cve_info['attack_complexity'] == 'HIGH':
            cve_info['attack_complexity'] = "Élevée"
        elif cve_info['attack_complexity'] == 'LOW':
            cve_info['attack_complexity'] = "Faible"
        
        if cve_info['privileges_required'] == 'HIGH':
            cve_info['privileges_required'] = "Élevés"
        elif cve_info['privileges_required'] == 'LOW':
            cve_info['privileges_required'] = 'Faible'
        elif cve_info['privileges_required'] == 'NONE':
            cve_info['privileges_required'] = 'Aucuns'

        if cve_info['user_interaction'] == 'NONE':
            cve_info['user_interaction'] = "Aucune"
        elif cve_info['user_interaction'] == 'REQUIRED':
            cve_info['user_interaction'] = 'Requise'

        if cve_info['scope'] == 'UNCHANGED':
            cve_info['scope'] = "Inchangé"
        elif cve_info['scope'] == 'CHANGED':
            cve_info['scope'] = "Modifié"

        if cve_info['confidentiality_impact'] == 'HIGH':
            cve_info['confidentiality_impact'] = "Élevée"
        elif cve_info['confidentiality_impact'] == 'LOW':
            cve_info['confidentiality_impact'] = 'Faible'
        elif cve_info['confidentiality_impact'] == 'NONE':
            cve_info['confidentiality_impact'] = 'Aucun'

        if cve_info['integrity_impact'] == 'HIGH':
            cve_info['integrity_impact'] = "Élevée"
        elif cve_info['integrity_impact'] == 'LOW':
            cve_info['integrity_impact'] = 'Faible'
        elif cve_info['integrity_impact'] == 'NONE':
            cve_info['integrity_impact'] = 'Aucun'
        
        if cve_info['availability_impact'] == 'HIGH':
            cve_info['availability_impact'] = "Élevée"
        elif cve_info['availability_impact'] == 'LOW':
            cve_info['availability_impact'] = 'Faible'
        elif cve_info['availability_impact'] == 'NONE':
            cve_info['availability_impact'] = 'Aucun'

        if cve_info['base_severity'] == 'HIGH':
            cve_info['base_severity'] = "Élevée"
        elif cve_info['base_severity'] == 'CRITICAL':
            cve_info['base_severity'] = 'Critique'

    return cve_list

#################################################################################
#                              PLAGE HORAIRE
#################################################################################

def plage():

    # Obtenir la date d'aujourd'hui
    aujourdhui = datetime.date.today()

    # Définir un délai (timedelta) de 1 jour
    un_jour = datetime.timedelta(days=1)

    # Obtenir la date d'hier
    hier = aujourdhui - un_jour

    # Extraire le jour en tant qu'entier
    # Obtenir la date d'aujourd'hui
    aujourdhui = datetime.datetime.now()

    # Obtenir la date d'hier
    hier = aujourdhui - datetime.timedelta(days=1)

    # Définir les heures, minutes et secondes pour les deux dates
    H = 9
    MIN = 0
    S = 0

    H2 = 8
    MIN2 = 59
    S2 = 59

    # Construire les objets datetime pour hier à 9h et aujourd'hui à 8h59
    start_datetime = datetime.datetime(hier.year, hier.month, hier.day, H, MIN, S)
    end_datetime = datetime.datetime(aujourdhui.year, aujourdhui.month, aujourdhui.day, H2, MIN2, S2)

    # Convertir les objets datetime en format de chaîne souhaité
    start_date = start_datetime.strftime('%Y-%m-%dT%H:%M:%S.000%z')
    end_date = end_datetime.strftime('%Y-%m-%dT%H:%M:%S.000%z')

    start_date2 = start_datetime.strftime("%Y%m%d")
    end_date2 = end_datetime.strftime("%Y%m%d")

    encoded_start_date = quote(start_date)
    encoded_end_date = quote(end_date)

    return encoded_start_date, encoded_end_date, start_date2, end_date2

def main():
    #definir les dates de début et de fin
    start_date, end_date, start_date2, end_date2 = plage()

    # récupérer les CVE dans un tableau
    cve_list = CVE(start_date, end_date)

    # mise en forme des valeurs
    cve_list = trad_vectors(cve_list)

    # création du powerpoint
    powerpoint(cve_list, start_date2, end_date2)

if __name__ == "__main__":
    main()
