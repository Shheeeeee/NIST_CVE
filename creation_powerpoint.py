from pptx import Presentation
from pptx.util import Pt

# Chemin d'accès au modèle PowerPoint
template_path = "Bulletin_de_veille_TEMPLATE.pptx"

# Crée une nouvelle présentation en utilisant le modèle
presentation = Presentation(template_path)


# Supposons que vous voulez modifier le titre de la deuxième diapositive (index 1)
slide_index = 1

# Accédez à la diapositive en utilisant l'index
slide = presentation.slides[slide_index]

# Trouvez la forme de titre de cette diapositive (la première forme de titre)

def titre():
    title_shape = slide.shapes.title

    # Modifiez le texte du titre
    nouveau_titre = "Nouveau titre de la diapo"
    title_shape.text = nouveau_titre

    # Définir le style du texte du titre
    font = title_shape.text_frame.paragraphs[0].font
    font.name = "Poppins SemiBold"
    #font.bold = True
    font.size = Pt(36)

    left = Pt(100)  # Définir la position horizontale (gauche) en points
    top = Pt(200)   # Définir la position verticale (haut) en points
    title_shape.left = left
    title_shape.top = top

# Enregistre la présentation modifiée dans un fichier
presentation.save("nouvelle_presentation.pptx")
