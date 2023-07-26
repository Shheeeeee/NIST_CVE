import requests
import datetime
from urllib.parse import quote
import time    

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches  # Import the Inches function
import copy

from googletrans import Translator

from itertools import cycle
from shutil import get_terminal_size
from threading import Thread
from time import sleep
from sys import stdout

#################################################################################
#                                 ANIMATION
#################################################################################

class Loader:
    def __init__(self, desc="Tentative de connexion ", end="", timeout=0.1):
        self.desc = desc
        self.end = end
        self.timeout = timeout

        self._thread = Thread(target=self._animate, daemon=True)
        self.steps = ["⢿", "⣻", "⣽", "⣾", "⣷", "⣯", "⣟", "⡿"]
        self.done = False

    def start(self):
        self._thread.start()
        return self

    def _animate(self):
        for c in cycle(self.steps):
            if self.done:
                break
            text = f"{self.desc}{self.attempt} {c}"
            print(f"\r{text}", flush=True, end="")
            stdout.flush()
            time.sleep(self.timeout)

    def __enter__(self):
        self.start()

    def stop(self):
        self.done = True
        cols = get_terminal_size((80, 20)).columns
        print("\r" + " " * cols, end="", flush=True)
        print(f"\r{self.end}", flush=True)

    def set_attempt(self, attempt):
        self.attempt = attempt

    def __exit__(self, exc_type, exc_value, tb):
        self.stop()

#################################################################################
#                                 API CALL
#################################################################################

def make_api_request(url, max_attempts=5):
    loader = Loader()
    loader_thread = Thread(target=loader.start, daemon=True)
    loader_thread.start()

    try:
        for attempt in range(1, max_attempts + 1):
            loader.set_attempt(attempt)
            #print(f"Connexion à l'API en cours...", end="")

            response = requests.get(url)

            if response.status_code == 200:
                data = response.json()
                loader.stop()
                print("\nConnexion réussie !")
                return data

            time.sleep(6)
    finally:
        loader.stop()
        if attempt == max_attempts:
            print(f"Impossible de se connecter à l'API après {max_attempts} tentatives. Veuillez réessayer plus tard.")
            exit()

#################################################################################
#                                    CVE
#################################################################################

def CVE(start_date, end_date):

    # Encoder les caractères spéciaux dans les dates


    # URL de l'API du NIST pour les CVE avec les dates de début et de fin spécifiées
    # url = f'https://services.nvd.nist.gov/rest/json/cves/2.0/?lastModStartDate={encoded_start_date}&lastModEndDate={encoded_end_date}'
    url = f'https://services.nvd.nist.gov/rest/json/cves/2.0/?pubStartDate={start_date}&pubEndDate={end_date}'

    # Nombre maximal de tentatives

    data = make_api_request(url)


    if len(data) != 0:
        cve_items = data['vulnerabilities']
        cve_list = []
        
        # Parcourir chaque CVE et extraire les informations
        for cve_item in cve_items:
            # Extract the severity (baseScore)
            cve_id = cve_item['cve']['id']

            severity = None
            for i in range(1, 100):
                key = f'cvssMetricV{i}'
                if key in cve_item['cve']['metrics']:
                    cve_metrics = cve_item['cve']['metrics'][key][0]

                    #################
                    # severity
                    #################
                    try:
                        severity = cve_metrics['baseScore']
                    except KeyError:
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                            severity = cve_metrics['cvssData']['baseScore']
                        except KeyError:
                            # Si 'attackVector' n'existe pas, définir attack_vector comme None
                            severity = None
                        continue

            # Print severity for debugging purposes
            if severity is None:
                print(f"Severity not found for CVE: {cve_item['cve']['id']}")


            # Skip processing the CVE if severity is less than 8
            if (severity is not None and severity < 8) or cve_item['cve']['vulnStatus'] == "Rejected" or cve_item['cve']['vulnStatus'] == "Modified":
                print(f"Skipping CVE {cve_id} due to low severity or rejected/modified status")
                continue

            cve_id = cve_item['cve']['id']
            vuln_status = cve_item['cve']['vulnStatus']

            # Initialize variables to store extracted information
            attack_vector = None
            vector_string = None
            descriptions = []
            
            
            # If vulnStatus is "Undergoing Analysis", et "severity" == None only extract and print descriptions
            if vuln_status == "Undergoing Analysis" and severity == None:
                if 'descriptions' in cve_item['cve']:
                    descriptions_data = cve_item['cve']['descriptions']
                    for description in descriptions_data:
                        descriptions.append(description['value'])

                print(f"CVE ID: {cve_id}")
                print(f"Vuln Status: {vuln_status}")
                print(f"Descriptions: {descriptions}")
                print()

                #################
                # il faut chercher des mots clefs dans la description tel que RCE, ...
                #################

            else:
                for i in range(1, 100): # Essayer différentes clés pour extraire le baseScore
                    key = f'cvssMetricV{i}'  # Construire la clé correspondante
                    if key in cve_item['cve']['metrics']:
                        cve_metrics = cve_item['cve']['metrics'][key][0]
                        
                        

                        # metrics = [cve_metrics['vectorString'], cve_metrics['attackVector'], cve_metrics['attackComplexity'], cve_metrics['privilegesRequired'], cve_metrics['userInteraction'], cve_metrics['scope'], cve_metrics['confidentialityImpact'], cve_metrics['integrityImpact'], cve_metrics['availabilityImpact'], cve_metrics['baseSeverity']]
                        # metrics2 = [cve_metrics['cvssData']['vectorString'], cve_metrics['cvssData']['attackVector'], cve_metrics['cvssData']['attackComplexity'], cve_metrics['cvssData']['privilegesRequired'], cve_metrics['cvssData']['userInteraction'], cve_metrics['cvssData']['scope'], cve_metrics['cvssData']['confidentialityImpact'], cve_metrics['cvssData']['integrityImpact'], cve_metrics['cvssData']['availabilityImpact'], cve_metrics['cvssData']['baseSeverity']]
                        # metrics3 = ['vectorString', 'attackVector', 'attackComplexity', 'privilegesRequired', 'userInteraction', 'scope', 'confidentialityImpact', 'integrityImpact', 'availabilityImpact', 'baseSeverity']
                        # for i in range(len(metrics)):
                        #     try:
                        #         # Essayer d'accéder à 'vectorString' dans cve_item['cve']['metrics'][key][0]['cvssData']
                        #         metrics3[i] = metrics[i]
                        #     except KeyError:
                        #         try:
                        #             # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                        #             metrics3[i] = metrics2[i]
                        #         except KeyError:
                        #             # Si 'attackVector' n'existe pas, définir attack_vector comme None
                        #             vector_string = None



                        #################
                        # vectorString
                        #################
                        try:
                            # Essayer d'accéder à 'vectorString' dans cve_item['cve']['metrics'][key][0]['cvssData']
                            vector_string = cve_metrics['vectorString']
                        except KeyError:
                            try:
                                # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                                vector_string = cve_metrics['cvssData']['vectorString']
                            except KeyError:
                                # Si 'attackVector' n'existe pas, définir attack_vector comme None
                                vector_string = None

                        #################
                        # attackVector
                        #################
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]
                            attack_vector = cve_metrics['attackVector']
                        except KeyError:
                            try:
                                # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                                attack_vector = cve_metrics['cvssData']['attackVector']
                            except KeyError:
                                # Si 'attackVector' n'existe pas, définir attack_vector comme None
                                attack_vector = None
                        
                        #################
                        # attack_complexity 
                        #################
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]
                            attack_complexity = cve_metrics['attackComplexity']
                        except KeyError:
                            try:
                                # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                                attack_complexity = cve_metrics['cvssData']['attackComplexity']
                            except KeyError:
                                # Si 'attackVector' n'existe pas, définir attack_vector comme None
                                attack_complexity = None
                        
                        #################
                        # privileges_required  
                        #################
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]
                            privileges_required = cve_metrics['privilegesRequired']
                        except KeyError:
                            try:
                                # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                                privileges_required = cve_metrics['cvssData']['privilegesRequired']
                            except KeyError:
                                # Si 'attackVector' n'existe pas, définir attack_vector comme None
                                privileges_required = None
                        
                        #################
                        # userInteraction  
                        #################
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]
                            user_interaction = cve_metrics['userInteraction']
                        except KeyError:
                            try:
                                # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                                user_interaction = cve_metrics['cvssData']['userInteraction']
                            except KeyError:
                                # Si 'attackVector' n'existe pas, définir attack_vector comme None
                                user_interaction = None
                        
                        #################
                        # scope   
                        #################
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]
                            scope = cve_metrics['scope']
                        except KeyError:
                            try:
                                # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                                scope = cve_metrics['cvssData']['scope']
                            except KeyError:
                                # Si 'attackVector' n'existe pas, définir attack_vector comme None
                                scope = None

                        #################
                        # confidentiality_impact   
                        #################
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]
                            confidentiality_impact = cve_metrics['confidentialityImpact']
                        except KeyError:
                            try:
                                # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                                confidentiality_impact = cve_metrics['cvssData']['confidentialityImpact']
                            except KeyError:
                                # Si 'attackVector' n'existe pas, définir attack_vector comme None
                                confidentiality_impact = None

                        #################
                        # integrityImpact   
                        #################
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]
                            integrity_impact = cve_metrics['integrityImpact']
                        except KeyError:
                            try:
                                # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                                integrity_impact = cve_metrics['cvssData']['integrityImpact']
                            except KeyError:
                                # Si 'attackVector' n'existe pas, définir attack_vector comme None
                                integrity_impact = None
                        
                        #################
                        # availability_impact    
                        #################
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]
                            availability_impact = cve_metrics['availabilityImpact']
                        except KeyError:
                            try:
                                # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                                availability_impact = cve_metrics['cvssData']['availabilityImpact']
                            except KeyError:
                                # Si 'attackVector' n'existe pas, définir attack_vector comme None
                                availability_impact = None

                        #################
                        # baseSeverity    
                        #################
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]
                            base_severity = cve_metrics['baseSeverity']
                        except KeyError:
                            try:
                                # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]['cvssData']
                                base_severity = cve_metrics['cvssData']['baseSeverity']
                            except KeyError:
                                # Si 'attackVector' n'existe pas, définir attack_vector comme None
                                base_severity = None

                        #################
                        # source    
                        #################
                        try:
                            # Essayer d'accéder à 'attackVector' dans cve_item['cve']['metrics'][key][0]
                            source = cve_item['cve']['references'][0]['url']
                        except KeyError:
                            source = None

                        
                        published = cve_item['cve']['published']
                        lastModified = cve_item['cve']['lastModified']

                # Extract 'descriptions'
                if 'descriptions' in cve_item['cve']:
                    descriptions = cve_item['cve']['descriptions'][0]['value']

                # Store all extracted information in a dictionary
                try:
                    cve_info = {
                    'cve_id': cve_id,
                    'published': published,
                    'lastModified': lastModified,
                    'vector_string': vector_string,
                    'attack_vector': attack_vector,
                    'attack_complexity': attack_complexity,
                    'privileges_required': privileges_required,
                    'user_interaction': user_interaction,
                    'scope': scope,
                    'confidentiality_impact': confidentiality_impact,
                    'integrity_impact': integrity_impact,
                    'availability_impact': availability_impact,
                    'severity': severity,
                    'base_severity': base_severity,
                    'descriptions': descriptions,
                    'source' : source
                    }

                    cve_list.append(cve_info)
                except Exception as e:
                    print("Une erreur s'est produite :", e)
                    continue
                    
        for cve_info in cve_list:
            print(f"CVE ID: {cve_info['cve_id']}")
            print(f"publiée : {cve_info['published']}  modifiée : {cve_info['lastModified']}")

            print(f"Vuln Status: {vuln_status}")
            print(f"Severity: {cve_info['severity']}")
            print(f"base_severity: {cve_info['base_severity']}")

            print(f"Vector String: {cve_info['vector_string']}")
            print(f"Attack Vector: {cve_info['attack_vector']}")
            print(f"attack_complexity: {cve_info['attack_complexity']}")
            print(f"privileges_required: {cve_info['privileges_required']}")
            print(f"user_interaction: {cve_info['user_interaction']}")
            print(f"scope: {cve_info['scope']}")
            print(f"confidentiality_impact: {cve_info['confidentiality_impact']}")
            print(f"integrity_impact: {cve_info['integrity_impact']}")
            print(f"availability_impact: {cve_info['availability_impact']}")

            print(f"Descriptions: {cve_info['descriptions']}")
            print(f"Source: {cve_info['source']}")
            print()

    
    else:
        print(f'Aucune CVE trouvée pour la plage de dates spécifiée.')
    
    return cve_list
    
#################################################################################
#                                 POWERPOINT
#################################################################################

def powerpoint(cve_list, start_date, end_date):
    
    # Le nom du template (modèle) et du fichier de sortie
    template_path = "Bulletin_de_veille_TEMPLATE.pptx"

    # Ouvrir le modèle
    presentation = Presentation(template_path)

    # Dupliquer la deuxième diapo (index 1) quatre fois
    for _ in range(len(cve_list) - 1):
        duplicate_slide(presentation, 1)

    diapo_index = 1  # Commencer à la deuxième diapo (index 1)

    for i in range(len(cve_list)):  # Modifier jusqu'à la cinquième diapo (index 4)
        slide = presentation.slides[diapo_index]
        
        cell_coords = [(2, 2), (3, 2), (4, 2), (5, 2), (2, 5), (3, 5), (4, 5), (5, 5), (7, 1), (10, 1)] 
        cell_values = [cve_list[diapo_index - 1]['attack_vector'], cve_list[diapo_index - 1]['attack_complexity'], cve_list[diapo_index - 1]['privileges_required'], cve_list[diapo_index - 1]['user_interaction'], cve_list[diapo_index - 1]['scope'], cve_list[diapo_index - 1]['confidentiality_impact'], cve_list[diapo_index - 1]['integrity_impact'], cve_list[diapo_index - 1]['availability_impact'], cve_list[diapo_index - 1]['descriptions'], cve_list[diapo_index - 1]['source']]  

        for i in range(len(cell_coords)):
            modify_table_cell_black(slide, cell_coords[i], cell_values[i])
            
        cell_coords = [(0, 1), (2, 7) , (9, 7)]  
        cell_values = [cve_list[diapo_index - 1]['vector_string'], f"{str(cve_list[diapo_index - 1]['severity'])}\n{cve_list[diapo_index - 1]['base_severity']}", cve_list[diapo_index - 1]['base_severity']]  

        for i in range(len(cell_coords)):
            modify_table_cell_white(slide, cell_coords[i], cell_values[i])
        
        titre(slide, cve_list[diapo_index - 1]['cve_id'])

        diapo_index += 1

    # Enregistrer la présentation modifiée dans un fichier
    nom_fichier = "CERT - Bulletin_de_veille_" + end_date + ".pptx"
    presentation.save(nom_fichier)
    print(f"La présentation a été sauvegardée dans '{nom_fichier}' avec succès.")

def titre(slide, titre):
    title_shape = slide.shapes.title

    # Modifiez le texte du titre

    title_shape.text = titre

    # Définir le style du texte du titre
    font = title_shape.text_frame.paragraphs[0].font
    font.name = "Poppins SemiBold"
    #font.bold = True
    font.size = Pt(36)

    left = Pt(55)  # Définir la position horizontale (gauche) en points
    top = Pt(90)   # Définir la position verticale (haut) en points
    title_shape.left = left
    title_shape.top = top

def duplicate_slide(pres, index):
    try:
        template = pres.slides[index]
        blank_slide_layout = template.slide_layout
    except:
        blank_slide_layout = pres.slide_layouts[0]

    copied_slide = pres.slides.add_slide(blank_slide_layout)

    for shp in template.shapes:
        if not shp.has_text_frame or shp.text_frame.text != "":  # Vérifier si la forme contient du texte
            el = shp.element
            newel = copy.deepcopy(el)
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return copied_slide


def modify_table_cell_black(slide, cell_coords, cell_value):
    row, col = cell_coords
    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break

    if table is None:
        print("Aucun tableau trouvé dans la diapositive.")
        return

    # Modifier le texte dans la cellule de la ligne et colonne spécifiées
    cell = table.cell(row, col)
    text_frame = cell.text_frame
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()

    # Modifier la mise en forme du texte
    run.text = str(cell_value)
    run.font.size = Pt(11)
    run.font.name = "Poppins"
    run.font.color.rgb = RGBColor(0, 0, 0)  # Noir


def modify_table_cell_white(slide, cell_coords, cell_value):
    row, col = cell_coords

    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break

    if table is None:
        print("Aucun tableau trouvé dans la diapositive.")
        return

    # Modifier le texte dans la cellule de la ligne et colonne spécifiées
    cell = table.cell(row, col)
    text_frame = cell.text_frame
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()

    # Modifier la mise en forme du texte
    run.text = str(cell_value)
    run.font.size = Pt(11)
    run.font.name = "Poppins"
    run.font.bold = True  # Pour mettre le texte en gras
    run.font.color.rgb = RGBColor(255, 255, 255)  # Pour mettre le texte en blanc (255, 255, 255 correspond au blanc en RGB)


#################################################################################
#                                 TRADUCTION
#################################################################################

def traduire_en_francais(texte):
    translator = Translator()
    try:
        traduction = translator.translate(texte, src='en', dest='fr')
        return traduction.text
    except:
        return texte

def traduire_donnees_en_francais(cve_list):
    for cve_info in cve_list:
        cve_info['base_severity'] = traduire_en_francais(str(cve_info['base_severity']))
        cve_info['attack_vector'] = traduire_en_francais(str(cve_info['attack_vector']))
        cve_info['attack_complexity'] = traduire_en_francais(str(cve_info['attack_complexity']))
        cve_info['privileges_required'] = traduire_en_francais(str(cve_info['privileges_required']))
        cve_info['user_interaction'] = traduire_en_francais(str(cve_info['user_interaction']))
        cve_info['scope'] = traduire_en_francais(str(cve_info['scope']))
        cve_info['confidentiality_impact'] = traduire_en_francais(str(cve_info['confidentiality_impact']))
        cve_info['integrity_impact'] = traduire_en_francais(str(cve_info['integrity_impact']))
        cve_info['availability_impact'] = traduire_en_francais(str(cve_info['availability_impact']))
        cve_info['descriptions'] = traduire_en_francais(str(cve_info['descriptions']))

    return cve_list

#################################################################################
#                              MISE EN FORME
#################################################################################

def mettre_majuscule_initiale(chaine):
    # Convertir la chaîne en minuscules
    chaine_minuscules = chaine.lower()
    # Mettre en majuscule la première lettre
    chaine_majuscule_initiale = chaine_minuscules.capitalize()
    return chaine_majuscule_initiale

def mettre_majuscule_initiale_tout(cve_list):
    for cve_info in cve_list:
        cve_info['base_severity'] = mettre_majuscule_initiale(str(cve_info['base_severity']))
        cve_info['attack_vector'] = mettre_majuscule_initiale(str(cve_info['attack_vector']))
        cve_info['attack_complexity'] = mettre_majuscule_initiale(str(cve_info['attack_complexity']))
        cve_info['privileges_required'] = mettre_majuscule_initiale(str(cve_info['privileges_required']))
        cve_info['user_interaction'] = mettre_majuscule_initiale(str(cve_info['user_interaction']))
        cve_info['scope'] = mettre_majuscule_initiale(str(cve_info['scope']))
        cve_info['confidentiality_impact'] = mettre_majuscule_initiale(str(cve_info['confidentiality_impact']))
        cve_info['integrity_impact'] = mettre_majuscule_initiale(str(cve_info['integrity_impact']))
        cve_info['availability_impact'] = mettre_majuscule_initiale(str(cve_info['availability_impact']))
        cve_info['descriptions'] = mettre_majuscule_initiale(str(cve_info['descriptions']))

    for cve_info in cve_list:
        if cve_info['base_severity'] == 'Haut':
            cve_info['base_severity'] = "Élevée"

        if cve_info['attack_complexity'] == 'Haut':
            cve_info['attack_complexity'] = "Élevée"
        
        if cve_info['privileges_required'] == 'Haut':
            cve_info['privileges_required'] = "Élevés"

        if cve_info['user_interaction'] == 'Aucun':
            cve_info['user_interaction'] = "Aucune"

        if cve_info['confidentiality_impact'] == 'Haut':
            cve_info['confidentiality_impact'] = "Élevée"

        if cve_info['integrity_impact'] == 'Haut':
            cve_info['integrity_impact'] = "Élevée"
        
        if cve_info['availability_impact'] == 'Haut':
            cve_info['availability_impact'] = "Élevée"

        if cve_info['scope'] == 'Unchanged':
            cve_info['scope'] = "Inchangé"

    return cve_list

#################################################################################
#                              PLAGE HORAIRE
#################################################################################

def plage():
    A = 2023
    M = 7
    J = 24
    H = 12
    MIN = 00
    S = 00

    A2 = 2023
    M2 = 7
    J2 = 25
    H2 = 12
    MIN2 = 00
    S2 = 00

    start_datetime = datetime.datetime(A, M, J, H, MIN, S)
    end_datetime = datetime.datetime(A2, M2, J2, H2, MIN2, S2)

    start_date = start_datetime.strftime('%Y-%m-%dT%H:%M:%S.000%z')
    end_date = end_datetime.strftime('%Y-%m-%dT%H:%M:%S.000%z')

    start_date2 = start_datetime.strftime("%Y%m%d")
    end_date2 = end_datetime.strftime("%Y%m%d")

    encoded_start_date = quote(start_date)
    encoded_end_date = quote(end_date)

    return encoded_start_date, encoded_end_date, start_date2, end_date2


def main():
    #definir les dates de début et de fin
    start_date, end_date, start_date2, end_date2 = plage()

    # récupérer les CVE dans un tableau
    cve_list = CVE(start_date, end_date)

    # Traduire tout les champs en anglais
    #cve_list = traduire_donnees_en_francais(cve_list)

    # mise en forme des valeurs
    cve_list = mettre_majuscule_initiale_tout(cve_list)

    # création du powerpoint
    powerpoint(cve_list, start_date2, end_date2)

if __name__ == "__main__":
    main()