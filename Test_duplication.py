from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import copy

def duplicate_slide(pres, index):
    try:
        template = pres.slides[index]
        blank_slide_layout = template.slide_layout
    except:
        blank_slide_layout = pres.slide_layouts[0]

    copied_slide = pres.slides.add_slide(blank_slide_layout)

    for shp in template.shapes:
        if not shp.has_text_frame or shp.text_frame.text != "":
            el = shp.element
            newel = copy.deepcopy(el)
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return copied_slide

def modify_table_cell(slide, row, col, cell_value):
    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break

    if table is None:
        print("Aucun tableau trouvé dans la diapositive.")
        return

    # Modifier le texte dans la cellule de la ligne et colonne spécifiées
    cell = table.cell(row, col)
    text_frame = cell.text_frame
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()

    # Modifier la mise en forme du texte
    run.text = cell_value
    run.font.size = Pt(11)
    run.font.name = "Poppins"
    run.font.color.rgb = RGBColor(0, 0, 0)  # Noir


def main():
    # Le nom du template (modèle) et du fichier de sortie
    template_path = "duplication.pptx"
    output_filename = "nouvelle_presentation.pptx"

    # Ouvrir le modèle
    presentation = Presentation(template_path)

    # Dupliquer la deuxième diapo (index 1) quatre fois
    for _ in range(4):
        duplicate_slide(presentation, 1)

    diapo_index = 2  # Commencer à la deuxième diapo (index 1)

    for i in range(5):  # Modifier jusqu'à la cinquième diapo (index 4)
        slide = presentation.slides[diapo_index]
        cell_value = f"Bonjour {i+1}"  # Par exemple, "Bonjour 1", "Bonjour 2", ...

        modify_table_cell(slide, 2, 2, cell_value)
        diapo_index += 1

    # Enregistrer la présentation modifiée dans un fichier
    presentation.save(output_filename)
    print(f"La présentation a été sauvegardée dans '{output_filename}' avec succès.")

if __name__ == "__main__":
    main()
